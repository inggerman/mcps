from __future__ import annotations

from pathlib import Path
from typing import Any

from mcp_shared.errors import ValidationError

_SUPPORTED = {".pdf", ".docx", ".pptx"}


def resolve_document(root: Path, path: str, max_file_size_mb: int) -> Path:
    base = root.resolve()
    document = (base / path).resolve() if not Path(path).is_absolute() else Path(path).resolve()
    if not document.is_relative_to(base):
        raise ValidationError(field="path", message="El documento está fuera de DOCUMENTS_ROOT.")
    if not document.is_file():
        raise ValidationError(field="path", message="El documento no existe.")
    if document.suffix.lower() not in _SUPPORTED:
        raise ValidationError(field="path", message="Formato soportado: PDF, DOCX o PPTX.")
    if document.stat().st_size > max_file_size_mb * 1024 * 1024:
        raise ValidationError(field="path", message="El documento supera el tamaño máximo.")
    return document


def extract_document(
    root: Path, path: str, max_file_size_mb: int, max_pages: int
) -> dict[str, Any]:
    document = resolve_document(root, path, max_file_size_mb)
    suffix = document.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(document, max_pages)
    if suffix == ".docx":
        return _extract_docx(document)
    return _extract_pptx(document, max_pages)


def _extract_pdf(path: Path, max_pages: int) -> dict[str, Any]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = [
        {"number": index + 1, "text": page.extract_text() or ""}
        for index, page in enumerate(reader.pages[:max_pages])
    ]
    return {
        "format": "pdf",
        "pages": pages,
        "page_count": len(reader.pages),
        "truncated": len(reader.pages) > max_pages,
        "text": "\n\n".join(str(page["text"]) for page in pages),
    }


def _extract_docx(path: Path) -> dict[str, Any]:
    from docx import Document

    document = Document(str(path))
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    tables = [
        [[cell.text for cell in row.cells] for row in table.rows] for table in document.tables
    ]
    return {
        "format": "docx",
        "paragraphs": paragraphs,
        "tables": tables,
        "text": "\n".join(paragraphs),
    }


def _extract_pptx(path: Path, max_pages: int) -> dict[str, Any]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    slides: list[dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides[:max_pages]):
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
        slides.append({"number": index + 1, "text": "\n".join(texts)})
    return {
        "format": "pptx",
        "slides": slides,
        "slide_count": len(presentation.slides),
        "truncated": len(presentation.slides) > max_pages,
        "text": "\n\n".join(slide["text"] for slide in slides),
    }


def get_document_metadata(root: Path, path: str, max_file_size_mb: int) -> dict[str, Any]:
    document = resolve_document(root, path, max_file_size_mb)
    stat = document.stat()
    return {
        "name": document.name,
        "format": document.suffix.lower().lstrip("."),
        "size": stat.st_size,
        "modified": stat.st_mtime,
    }


# ---------------------------------------------------------------------------
# Tools nuevos
# ---------------------------------------------------------------------------


def list_documents(root: Path, directory: str = "") -> list[dict[str, Any]]:
    """Lista documentos soportados en un directorio."""
    base = root.resolve()
    target = (base / directory).resolve() if directory else base
    if not target.is_relative_to(base):
        raise ValidationError(field="directory", message="Directorio fuera de DOCUMENTS_ROOT.")
    if not target.is_dir():
        raise ValidationError(field="directory", message="Directorio no existe.")

    documents: list[dict[str, Any]] = []
    for f in target.rglob("*"):
        if f.is_file() and f.suffix.lower() in _SUPPORTED:
            stat = f.stat()
            documents.append({
                "name": f.name,
                "path": str(f.relative_to(base)),
                "format": f.suffix.lower().lstrip("."),
                "size": stat.st_size,
                "size_mb": round(stat.st_size / (1024 * 1024), 2),
                "modified": stat.st_mtime,
            })
    return documents


def search_in_document(root: Path, path: str, query: str, max_file_size_mb: int) -> dict[str, Any]:
    """Busca texto dentro de un documento."""
    extracted = extract_document(root, path, max_file_size_mb, 5000)
    text = extracted.get("text", "")
    query_lower = query.lower()
    text_lower = text.lower()

    matches: list[dict[str, Any]] = []
    lines = text.split("\n")
    for i, line in enumerate(lines, 1):
        if query_lower in line.lower():
            matches.append({"line": i, "text": line.strip()[:200]})

    return {
        "file": path,
        "query": query,
        "total_matches": len(matches),
        "matches": matches[:50],
    }


def count_pages(root: Path, path: str, max_file_size_mb: int) -> dict[str, Any]:
    """Cuenta paginas o slides de un documento."""
    document = resolve_document(root, path, max_file_size_mb)
    suffix = document.suffix.lower()

    if suffix == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(str(document))
        return {"file": path, "format": "pdf", "count": len(reader.pages)}
    elif suffix == ".docx":
        from docx import Document
        doc = Document(str(document))
        return {"file": path, "format": "docx", "paragraphs": len(doc.paragraphs), "tables": len(doc.tables)}
    elif suffix == ".pptx":
        from pptx import Presentation
        pres = Presentation(str(document))
        return {"file": path, "format": "pptx", "slides": len(pres.slides)}
    return {"file": path, "error": "Unknown format"}


def extract_text_only(root: Path, path: str, max_file_size_mb: int, max_pages: int) -> str:
    """Extrae solo el texto plano de un documento."""
    result = extract_document(root, path, max_file_size_mb, max_pages)
    return result.get("text", "")


def get_document_summary(root: Path, path: str, max_file_size_mb: int, max_pages: int) -> dict[str, Any]:
    """Genera un resumen basico de un documento."""
    result = extract_document(root, path, max_file_size_mb, max_pages)
    text = result.get("text", "")
    words = text.split()
    lines = text.split("\n")

    return {
        "file": path,
        "format": result.get("format"),
        "total_words": len(words),
        "total_lines": len(lines),
        "total_characters": len(text),
        "page_count": result.get("page_count", result.get("slide_count", 0)),
        "truncated": result.get("truncated", False),
        "first_500_chars": text[:500],
    }


def convert_to_markdown(root: Path, path: str, max_file_size_mb: int, max_pages: int) -> str:
    """Convierte un documento a formato Markdown."""
    result = extract_document(root, path, max_file_size_mb, max_pages)
    fmt = result.get("format")

    if fmt == "pdf":
        lines = [f"# {path}\n"]
        for page in result.get("pages", []):
            lines.append(f"## Page {page['number']}\n")
            lines.append(page["text"])
            lines.append("")
        return "\n".join(lines)
    elif fmt == "docx":
        lines = [f"# {path}\n"]
        for para in result.get("paragraphs", []):
            lines.append(para)
            lines.append("")
        for i, table in enumerate(result.get("tables", []), 1):
            lines.append(f"### Table {i}\n")
            for row in table:
                lines.append("| " + " | ".join(row) + " |")
            lines.append("")
        return "\n".join(lines)
    elif fmt == "pptx":
        lines = [f"# {path}\n"]
        for slide in result.get("slides", []):
            lines.append(f"## Slide {slide['number']}\n")
            lines.append(slide["text"])
            lines.append("")
        return "\n".join(lines)
    return f"# {path}\n\nUnsupported format"


def batch_extract(root: Path, directory: str, max_file_size_mb: int, max_pages: int) -> dict[str, Any]:
    """Extrae texto de todos los documentos en un directorio."""
    docs = list_documents(root, directory)
    results: list[dict[str, Any]] = []

    for doc in docs:
        try:
            summary = get_document_summary(root, doc["path"], max_file_size_mb, max_pages)
            results.append({
                "file": doc["path"],
                "format": doc["format"],
                "words": summary["total_words"],
                "pages": summary["page_count"],
            })
        except Exception as exc:
            results.append({"file": doc["path"], "error": str(exc)[:100]})

    return {
        "total_documents": len(docs),
        "processed": len(results),
        "results": results,
    }


def get_document_stats(root: Path, directory: str = "") -> dict[str, Any]:
    """Retorna estadisticas de los documentos en un directorio."""
    docs = list_documents(root, directory)

    by_format: dict[str, int] = {}
    total_size = 0
    for doc in docs:
        fmt = doc["format"]
        by_format[fmt] = by_format.get(fmt, 0) + 1
        total_size += doc["size"]

    return {
        "total_documents": len(docs),
        "by_format": by_format,
        "total_size_bytes": total_size,
        "total_size_mb": round(total_size / (1024 * 1024), 2),
        "average_size_mb": round(total_size / max(len(docs), 1) / (1024 * 1024), 2),
    }


def validate_document(root: Path, path: str, max_file_size_mb: int) -> dict[str, Any]:
    """Valida que un documento sea accesible y soportado."""
    document = resolve_document(root, path, max_file_size_mb)
    stat = document.stat()

    issues: list[str] = []
    if stat.st_size == 0:
        issues.append("File is empty.")
    if document.suffix.lower() not in _SUPPORTED:
        issues.append(f"Unsupported format: {document.suffix}")

    return {
        "file": path,
        "valid": len(issues) == 0,
        "issues": issues,
        "format": document.suffix.lower().lstrip("."),
        "size_bytes": stat.st_size,
    }


def extract_tables(root: Path, path: str, max_file_size_mb: int) -> dict[str, Any]:
    """Extrae tablas de un documento (DOCX principalmente)."""
    document = resolve_document(root, path, max_file_size_mb)
    if document.suffix.lower() != ".docx":
        return {"file": path, "tables": [], "message": "Tables extraction only supported for DOCX."}

    from docx import Document
    doc = Document(str(document))
    tables = []
    for i, table in enumerate(doc.tables, 1):
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        tables.append({"table_number": i, "rows": len(rows), "data": rows})

    return {"file": path, "tables_count": len(tables), "tables": tables}


def extract_images_info(root: Path, path: str, max_file_size_mb: int) -> dict[str, Any]:
    """Retorna informacion sobre imagenes en un documento."""
    document = resolve_document(root, path, max_file_size_mb)
    suffix = document.suffix.lower()

    if suffix == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(str(document))
        image_count = 0
        for page in reader.pages:
            if "/XObject" in page.get("/Resources", {}):
                x_objects = page["/Resources"]["/XObject"].get_object()
                for obj in x_objects:
                    if x_objects[obj].get("/Subtype") == "/Image":
                        image_count += 1
        return {"file": path, "format": "pdf", "images_found": image_count}
    elif suffix == ".docx":
        from docx import Document
        doc = Document(str(document))
        image_count = 0
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image_count += 1
        return {"file": path, "format": "docx", "images_found": image_count}
    elif suffix == ".pptx":
        from pptx import Presentation
        pres = Presentation(str(document))
        image_count = 0
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    image_count += 1
        return {"file": path, "format": "pptx", "images_found": image_count}
    return {"file": path, "images_found": 0}


def compare_documents(root: Path, path_a: str, path_b: str, max_file_size_mb: int, max_pages: int) -> dict[str, Any]:
    """Compara dos documentos."""
    result_a = extract_document(root, path_a, max_file_size_mb, max_pages)
    result_b = extract_document(root, path_b, max_file_size_mb, max_pages)

    text_a = result_a.get("text", "")
    text_b = result_b.get("text", "")

    words_a = set(text_a.split())
    words_b = set(text_b.split())

    common = words_a & words_b
    unique_a = words_a - words_b
    unique_b = words_b - words_a

    similarity = len(common) / max(len(words_a | words_b), 1) * 100

    return {
        "file_a": path_a,
        "file_b": path_b,
        "format_a": result_a.get("format"),
        "format_b": result_b.get("format"),
        "words_a": len(words_a),
        "words_b": len(words_b),
        "common_words": len(common),
        "unique_to_a": len(unique_a),
        "unique_to_b": len(unique_b),
        "similarity_percent": round(similarity, 2),
    }


def export_document_text(root: Path, path: str, max_file_size_mb: int, max_pages: int) -> dict[str, Any]:
    """Exporta el texto de un documento en formato estructurado."""
    result = extract_document(root, path, max_file_size_mb, max_pages)
    return {
        "file": path,
        "format": result.get("format"),
        "text": result.get("text", ""),
        "metadata": {
            "pages": result.get("page_count", result.get("slide_count", 0)),
            "truncated": result.get("truncated", False),
            "paragraphs": len(result.get("paragraphs", [])),
            "tables": len(result.get("tables", [])),
        },
    }
